'''
Generate structure breakdown cards for each model based on prediction results and experimental data.
Note for some models we had replaced some standardised smiles with XX in the input to make the predictions fail.
'''
# setup logging
import logger
import logging
log = logger.setup_applevel_logger(logger_name = '02_structure_breakdown', file_name ='logs/02_structure_breakdown.log', level_stream=logging.INFO, level_file=logging.DEBUG)

import pandas as pd
import numpy as np

from rdkit import Chem
from rdkit.Chem import Descriptors

from pptx import Presentation
from pdf2image import convert_from_path

import subprocess
import os

# pandas display options
# do not fold dataframes
pd.set_option('expand_frame_repr',False)
# maximum number of columns
pd.set_option("display.max_columns",50)
# maximum number of rows
pd.set_option("display.max_rows",500)
# precision of float numbers
pd.set_option("display.precision",3)
# maximum column width
pd.set_option("max_colwidth", 250)

# enable pandas copy-on-write
pd.options.mode.copy_on_write = True



def replace_placeholders_in_runs(shape, values):
    """
    Replace placeholders inside a shape without changing formatting.
    """
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        # get the full paragraph text (merged across runs)
        full_text = "".join(run.text for run in paragraph.runs)

        # replace all placeholders
        for key, val in values.items():
            if key in full_text:
                full_text = full_text.replace(key, val)

        # now rewrite runs carefully:
        # delete all existing runs
        for run in paragraph.runs:
            run.text = ""

        # insert the new text into the FIRST run to preserve formatting
        paragraph.runs[0].text = full_text

def iter_shapes(shapes):
    """
    Recursively iterate through shapes, including group shapes.
    """
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # group
            for subshape in iter_shapes(shape.shapes):
                yield subshape


# convert pptx to pdf using libre office
def pptx_to_pdf_libreoffice(input_file, output_dir=None):
    libreoffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"

    if output_dir is None:
        output_dir = os.path.dirname(input_file)

    subprocess.run([
        libreoffice_path,  # or "libreoffice" depending on your installation
        "--headless",
        "--convert-to", "pdf",
        input_file,
        "--outdir", output_dir
    ], check=True)

    log.info(f"Converted {input_file} to PDF in {output_dir}")


# read in all structures for which we generated predictions
smiles = pd.read_excel(r'data/structures/smiles.xlsx')


# read all predictions
prediction_files = [r'data/predictions/vega/processed/predictions_vega.xlsx',
                    r'data/predictions/ecosar/processed/predictions_ecosar.xlsx',
                    r'data/predictions/isaferat/processed/predictions_isaferat_AD_no_extrapolation.xlsx',
                    r'data/predictions/isaferat/processed/predictions_isaferat_AD_extrapolation_all.xlsx',
                    r'data/predictions/isaferat/processed/predictions_isaferat_AD_extrapolation_structural_in_domain.xlsx',
                    r'data/predictions/test/relax_fragment_constraint_false/processed/predictions_test.xlsx',
                    # r'data/predictions/test/relax_fragment_constraint_true/processed/predictions_test.xlsx',
                    r'data/predictions/kate/processed/predictions_kate.xlsx',
                    'data/predictions/trident/processed/predictions_trident.xlsx',]
prediction_data = []
for prediction_file in prediction_files:
    prediction_data.append(pd.read_excel(prediction_file))
prediction_data = pd.concat(prediction_data, axis='index', ignore_index=True, sort=False)
# keep only the numerical predictions
msk = prediction_data['predicted quantity'].str.contains(r'\((?:m|u)g/L\)')
prediction_data = prediction_data.loc[msk]
# keep the most conservative prediction if multiple predictions are available for the same standardised SMILES
prediction_data = (prediction_data
                   .sort_values(['platform', 'model name', 'smiles (standardised)', 'study type', 'prediction'], ascending=[True, True, True, True, True])
                   .drop_duplicates(subset=['platform', 'model name', 'smiles (standardised)', 'study type', ], keep='first')
                   .reset_index(drop=True)
                   )

# read the acute experimental data
experimental_data_acute = pd.read_pickle(r'data/fish_acute/processed/REACH_short_term_fish_measurement.pickle')
# keep only quantitative AFT data
msk = experimental_data_acute['scenario'].isin(['AFT_1a', 'AFT_1b']) & experimental_data_acute['exclusion reasons'].isnull()
# take the most conservative value if multiple measurements are available for the same standardised SMILES
experimental_data_acute = (experimental_data_acute
                           .loc[msk]
                           .groupby(['smiles'], dropna=False)
                           [['effect concentration (mol/L, lower bound)', 'CAS number']]
                           .agg(**{'effect concentration (mol/L)': pd.NamedAgg(column='effect concentration (mol/L, lower bound)', aggfunc='min'),
                                   'CAS number': pd.NamedAgg(column='CAS number', aggfunc=lambda x: ', '.join(sorted(set(x.dropna()))))}
                                )
                           .reset_index()
                           # keep only the DSSTox smiles that could be standardised
                           .merge(smiles[['smiles', 'smiles (standardised)']], how='inner', left_on='smiles', right_on='smiles')
                           .assign(**{'molecular weight (standardised)': lambda df: df['smiles (standardised)'].apply(lambda x: Chem.MolFromSmiles(x)).apply(lambda x: Descriptors.MolWt(x) if x is not None else np.nan)})
                           #  compute the effect levels in mg/L for the standardised structures
                           .assign(**{'effect concentration (mg/L)': lambda df: df['effect concentration (mol/L)']*df['molecular weight (standardised)']*1000})
                           .drop('smiles', axis='columns')
                           )
# read the chronic experimental data
experimental_data_chronic = pd.read_pickle(r'data/fish_chronic/processed/REACH_long_term_fish_measurement.pickle')
# keep only quantitative FELS data
msk = experimental_data_chronic['scenario'].isin(['FLS_1a', 'FLS_1b']) & experimental_data_chronic['exclusion reasons'].isnull()
# take the most conservative value if multiple measurements are available for the same standardised SMILES
experimental_data_chronic = (experimental_data_chronic
                           .loc[msk]
                           .groupby(['smiles'], dropna=False)
                           [['effect concentration (mol/L, lower bound)', 'CAS number']]
                           .agg(**{'effect concentration (mol/L)': pd.NamedAgg(column='effect concentration (mol/L, lower bound)', aggfunc='min'),
                                   'CAS number': pd.NamedAgg(column='CAS number', aggfunc=lambda x: ', '.join(sorted(set(x.dropna()))))}
                                )
                           .reset_index()
                           # keep only the DSSTox smiles that could be standardised
                           .merge(smiles[['smiles', 'smiles (standardised)']], how='inner', left_on='smiles', right_on='smiles')
                           .assign(**{'molecular weight (standardised)': lambda df: df['smiles (standardised)'].apply(lambda x: Chem.MolFromSmiles(x)).apply(lambda x: Descriptors.MolWt(x) if x is not None else np.nan)})
                           #  compute the effect levels in mg/L for the standardised structures
                           .assign(**{'effect concentration (mg/L)': lambda df: df['effect concentration (mol/L)'] * df['molecular weight (standardised)']*1000})
                           .drop('smiles', axis='columns')
                           )

# put all experimental results together
experimental_data = pd.concat([experimental_data_acute.assign(**{'study type': 'acute'}),
                                    experimental_data_chronic.assign(**{'study type': 'chronic'})], axis='index', ignore_index=True, sort=False)


models = prediction_data[['platform', 'model name', 'study type']].drop_duplicates().sort_values(by=['platform', 'model name']).reset_index(drop=True)

# generate the structure breakdown card for each model
all_values = []
for idx, row in models.iterrows():
    log.info('processing model %d/%d: %s - %s', idx+1, len(models), row['platform'], row['model name'])
    # initialize the values dictionary
    values = dict()

    platform = row['platform']
    model_name = row['model name']
    study_type = row['study type']
    values['{{platform}}'] = platform
    values['{{model}}'] = model_name
    values['{{study_type}}'] = study_type
    # get the predictions for this model
    msk = (prediction_data['platform'] == platform) & (prediction_data['model name'] == model_name)
    prediction_data_filtered = prediction_data.loc[msk]

    # convert units if necessary
    if prediction_data_filtered['predicted quantity'].iloc[0] in ['LC50 (mg/L)', 'ChV (mg/L)', 'NOEC (mg/L)', 'EC10 (mg/L)']:
        pass
    elif prediction_data_filtered['predicted quantity'].iloc[0] == 'AC50 (ug/L)':
        prediction_data_filtered['prediction'] = prediction_data_filtered['prediction']*1.e-3
    else:
        raise ValueError(f'Unknown predicted quantity {prediction_data_filtered["predicted quantity"].iloc[0]} for model {model_name}')


    # number and percentage of structures that succeeded and failed
    n = len(prediction_data_filtered)
    n_f = len(prediction_data_filtered.loc[prediction_data_filtered['prediction status'] == 'failed'])
    n_s = len(prediction_data_filtered.loc[prediction_data_filtered['prediction status'] == 'succeeded'])
    # percentage of failed predictions in round number
    p_f = f'{100*n_f/n:.0f}'
    p_s = f'{100*n_s/n:.0f}'
    values.update({'{{n}}': str(n),
                   '{{n_f}}': str(n_f),
                   '{{n_s}}': str(n_s),
                   '{{p_f}}': p_f,
                   '{{p_s}}': p_s})

    # keep only the successful predictions
    msk = prediction_data_filtered['prediction status'] == 'succeeded'
    prediction_data_filtered = prediction_data_filtered.loc[msk]

    # set the scenario
    values['{{scenario}}'] = 'AFT 1a/b' if study_type == 'acute' else 'FLS 1a/b'

    # number of structures with predictions and experimental data
    msk = (experimental_data['study type'] == study_type)
    experimental_data_filtered = experimental_data.loc[msk]
    n_sna = prediction_data_filtered.merge(experimental_data_filtered, on=['smiles (standardised)', 'study type'], how='left', indicator=True).query('_merge=="left_only"').shape[0]
    n_sa = prediction_data_filtered.merge(experimental_data_filtered, on=['smiles (standardised)', 'study type'], how='left', indicator=True).query('_merge=="both"').shape[0]
    values['{{n_sna}}'] = str(n_sna)
    values['{{n_sa}}'] = str(n_sa)

    # keep only the predictions for which experimental data is available
    predictions_experimental_filtered = prediction_data_filtered.merge(experimental_data_filtered, on='smiles (standardised)', how='inner')

    # predictions with no effect at saturation
    msk = predictions_experimental_filtered['no effects at saturation'] == 'yes'
    n_nes = msk.sum()
    p_nes = f'{100*n_nes/len(predictions_experimental_filtered):.0f}'
    values.update({'{{n_nes}}': str(n_nes),
                   '{{p_nes}}': p_nes})

    # keep only the predictions that are not "no effects at saturation"
    n_es = len(predictions_experimental_filtered) - n_nes
    p_es = f'{100*n_es/len(predictions_experimental_filtered):.0f}'
    predictions_experimental_filtered = predictions_experimental_filtered.loc[~msk]
    values.update({'{{n_es}}': str(n_es),
                   '{{p_es}}': p_es})

    # predictions in training set
    msk = predictions_experimental_filtered['training/validation set'] == 'training set'
    n_t = msk.sum()
    values.update({'{{n_t}}': str(n_t)})

    # predictions in validation set
    msk = predictions_experimental_filtered['training/validation set'] == 'validation set'
    n_v = msk.sum()
    values.update({'{{n_v}}': str(n_v)})

    # keep only predictions not in training or validation set
    msk = predictions_experimental_filtered['training/validation set'] == 'not in training/validation set'
    predictions_experimental_filtered = predictions_experimental_filtered.loc[msk]
    n_ntv = len(predictions_experimental_filtered)
    values.update({'{{n_ntv}}': str(n_ntv)})

    # out of domain predictions
    msk = predictions_experimental_filtered['AD'] == 'out of domain'
    n_od = msk.sum()
    p_od = f'{100*n_od/n_ntv:.0f}'
    n_id = n_ntv - n_od
    p_id = f'{100*n_id/n_ntv:.0f}'
    values.update({'{{n_od}}': str(n_od),
                   '{{p_od}}': p_od,
                   '{{n_id}}': str(n_id),
                   '{{p_id}}': p_id})

    # keep only the predictions that are in domain
    msk = predictions_experimental_filtered['AD'] == 'in domain'
    predictions_experimental_filtered = predictions_experimental_filtered.loc[msk]
    predictions_experimental_filtered['prediction'] = pd.to_numeric(predictions_experimental_filtered['prediction'], errors='raise')
    predictions_experimental_filtered['effect concentration (mg/L)'] = pd.to_numeric(predictions_experimental_filtered['effect concentration (mg/L)'], errors='raise')

    if not predictions_experimental_filtered.empty:
        # predictions within 1 log unit from experimental value
        msk = (np.abs(np.log10(predictions_experimental_filtered['prediction'])
                      - np.log10(predictions_experimental_filtered['effect concentration (mg/L)'])) <= 1)
        n_1lu = msk.sum()
        p_1lu = f'{100*n_1lu/len(predictions_experimental_filtered):.0f}'
        values.update({'{{n_1lu}}': str(n_1lu),
                       '{{p_1lu}}': p_1lu})

        # predictions that over predict toxicity by more than 1 log unit
        msk = (np.log10(predictions_experimental_filtered['prediction'])
                  < np.log10(predictions_experimental_filtered['effect concentration (mg/L)']) - 1)
        n_op = msk.sum()
        p_op = f'{100*n_op/len(predictions_experimental_filtered):.0f}'
        values.update({'{{n_op}}': str(n_op),
                       '{{p_op}}': p_op})

        # predictions that under predict toxicity by more than 1 log unit
        msk = (np.log10(predictions_experimental_filtered['prediction'])
                    > np.log10(predictions_experimental_filtered['effect concentration (mg/L)']) + 1)
        n_up = msk.sum()
        p_up = f'{100*n_up/len(predictions_experimental_filtered):.0f}'
        values.update({'{{n_up}}': str(n_up),
                    '{{p_up}}': p_up})
    else:
        values.update({'{{n_1lu}}': '-',
                       '{{p_1lu}}': '-',
                       '{{n_op}}': '-',
                       '{{p_op}}': '-',
                       '{{n_up}}': '-',
                       '{{p_up}}': '-'})

    # load the template
    prs = Presentation("figures/structure_breakdown/structure_breakdown_template.pptx")
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            replace_placeholders_in_runs(shape, values)

    # save the output as pptx
    prs.save(fr'figures/structure_breakdown/{platform}_{model_name}.pptx')

    # convert the pptx to pdf
    pptx_to_pdf_libreoffice(fr'figures/structure_breakdown/{platform}_{model_name}.pptx')

    # convert the pdf to png (assuming that the pdf has only one page)
    log.info('converting pdf to png for model %s - %s', platform, model_name)
    pdf_path = fr'figures/structure_breakdown/{platform}_{model_name}.pdf'
    images = convert_from_path(pdf_path, dpi=600, poppler_path=r'D:/Applications/poppler/Release-25.11.0-0/Library/bin')
    for i, image in enumerate(images):
        image.save(fr'figures/structure_breakdown/{platform}_{model_name}.png', "PNG", use_temp=False)
        break
    all_values.append(values)

all_values = pd.DataFrame(all_values)
all_values.columns = [col.replace('{','').replace('}','') for col in all_values.columns]
all_values['AD coverage'] = (all_values['n_s'].astype(int)/all_values['n'].astype(int))*(all_values['n_id'].astype(int)/all_values['n_ntv'].astype(int))
all_values.to_excel(r'data/predictions/structure_breakdown_summary.xlsx', index=False)